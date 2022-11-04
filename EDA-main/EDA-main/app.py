# -*- coding: utf-8 -*-
import streamlit_authenticator as stauth
import streamlit as st
import pandas as pd
import numpy as np
import seaborn as sns
from matplotlib import pyplot as plt
import nbformat as nbf
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import database as db

# -----User Authentication ---------

users =db.fetch_all_users()

usernames = [user["key"] for user in users]
names = [user["name"] for user in users]
hashed_passwords = [user["password"] for user in users]

authenticator = stauth.Authenticate(names, usernames, hashed_passwords, "EDA", "abcdef", cookie_expiry_days=30)

name, authentication_status, username = authenticator.login("Login","main")
if authentication_status == False:
    st.error("Username/Password is incorrect")

if authentication_status == None:
    st.warning("Please enter your Username and Password")

if authentication_status == True:
    hide_st_style = """
                    <style>
                    #MainMenu{visibility: hidden; }
                    footer{visibility: hidden;}
                    </style>
                    """
    st.markdown(hide_st_style, unsafe_allow_html= True)
    st.set_option('deprecation.showfileUploaderEncoding', False)
    st.set_option('deprecation.showPyplotGlobalUse', False)
    @st.cache

    def get_data(filename):
        """
        Function to extract CSV filename and convert to pandas dataframe.

        :param filename: IOString Path to CSV path.
        :return: {pd.DataFrame}:  Pandas Dataframe.
        """
        return pd.read_csv(filename)

    def main():
        """
        Main function for a Universal EDA App. Run this to run the app.
        """

        st.title('📊WEBEDA')
        st.subheader('Web based tool for general purpose Exploratory Data Analysis')

        uploaded_file = st.file_uploader('Upload CSV file to begin (Max file size allowed: 200MB)', type='csv')

        if uploaded_file is not None:
            df = get_data(uploaded_file)
            
            authenticator.logout("Logout", "sidebar")
            st. sidebar.title(f"Welcome {name}")
            st.sidebar.title('Tools 🔧')
            if st.checkbox('Show raw data', value=False):
                st.write(df)

            target_column = st.selectbox('Select Target Column', list(df.columns), key='target_column')
            if target_column is not None:
                if st.sidebar.checkbox('✍ Describe', value=False):
                    st.markdown('## Data Description')
                    st.write(df.describe())
                    st.markdown('### Columns that are potential binary features')
                    bin_cols = []
                    for col in df.columns:
                        if len(df[col].value_counts()) == 2:
                            bin_cols.append(col)
                    st.write(bin_cols)
                    st.markdown('### Columns Types')
                    st.write(df.astype(str))

                if st.sidebar.checkbox('👁 Missing Data', value=False):
                    st.markdown('## Missing Data')
                    total = df.isnull().sum().sort_values(ascending=False)
                    percent = (df.isnull().sum() / df.isnull().count()).sort_values(ascending=False)
                    missing_data = pd.concat([total, percent], axis=1, keys=['Total', 'Percent'])
                    st.write(missing_data)

                if st.sidebar.checkbox('🔢 Value Counts', value=False):
                    st.markdown('## Value Counts')
                    col = st.selectbox('Select Column', list(df.columns), key= None)
                    st.write(df[col].value_counts())

                if st.sidebar.checkbox('🧬 Unique Elements', value=False):
                    st.markdown('## Unique elements')
                    if st.checkbox('Show all unique elements', value=False):
                        st.write(df.nunique())
                    col = st.selectbox('Show columnwise unique elements', list(df.columns), key='unique_col')
                    st.write(df[col].unique())

                if st.sidebar.checkbox('〽 Show Distribution', False):
                    st.subheader(f'Distribution of {target_column}')
                    try:
                        sns.distplot(df[target_column])
                        st.write("Skewness: %.3f" % df[target_column].skew())
                        st.write("Kurtosis: %.3f" % df[target_column].kurt())
                        st.pyplot()
                    except:
                        st.error('Invalid Column')

                if st.sidebar.checkbox('📈 Scatter Plot', value=False):
                    st.markdown('## Scatter Plots')
                    scatter_cols = st.multiselect('Select Column', list(df.columns), key='scatter_cols')
                    for col in scatter_cols:
                        try:
                            data = pd.concat([df[target_column], df[col]], axis=1)
                            data.plot.scatter(x=col, y=target_column)
                            st.pyplot()
                        except:
                            st.error('Invalid column')

                if st.sidebar.checkbox('🈁 Box Plot', value=False):
                    st.markdown('## Box Plots')
                    box_cols = st.multiselect('Select Column', list(df.columns), key='box_cols')
                    for col in box_cols:
                        try:
                            data = pd.concat([df[target_column], df[col]], axis=1)
                            f, ax = plt.subplots(figsize=(8, 6))
                            fig = sns.boxplot(x=col, y=target_column, data=data)
                            fig.axis(ymin=np.min(df[target_column]), ymax=np.max(df[target_column]))
                            st.pyplot()
                        except:
                            st.error('Invalid column')

                if st.sidebar.checkbox('➿ Pair Plot', value=False):
                    st.markdown('## Pair Plots')
                    pair_cols = st.multiselect('Select Column', list(df.columns), key='pair_plot')
                    plot_size = st.number_input('Select Plot size', 1.0, 5.0, step=0.5, key='plot_size', value=2.5)
                    cols = [target_column]
                    for col in pair_cols:
                        cols.append(col)
                    try:
                        sns.set()
                        sns.pairplot(df[cols], size=plot_size)
                        st.pyplot()
                    except:
                        st.error('Invalid column')

                if st.sidebar.checkbox('🧮 Correlation Matrix', value=False):
                    st.markdown('## Correlation matrix (heatmap style)')
                    eda_cols = list(df.columns)
                    if len(eda_cols) > 5:
                        eda_cols = eda_cols[:4]
                    eda_cols = st.multiselect("Select features:", df.columns.tolist(), default=eda_cols)
                    plt.figure(figsize=(12, 8))
                    ax = sns.heatmap(df[eda_cols].corr(), cmap="OrRd", linecolor='white', linewidths=1)
                    st.pyplot()

                if st.sidebar.checkbox('🖼️ Create PPT', value=False):
                    st.markdown('### Create a Power Point Presentation\n'
                                'Which EDA Tools you want to include in the PPTX?')

                    prs = Presentation()
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    subtitle = slide.placeholders[1]
                    title.text = "EDA Presentation"
                    subtitle.text = "Credits to Aditya, Nikhil, Shantanu"
                    top = Inches(7)
                    left = Inches(4.5)
                    width = height = Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = "<Project Number> <Project name>\n<Reference number> <Issue> <Last updated>"
                    tf.paragraphs[0].font.size = Pt(8)
                    tf.paragraphs[1].font.size = Pt(8)
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    tf.paragraphs[1].alignment = PP_ALIGN.CENTER

                    if st.checkbox('📈: Scatter Plot', value=False):

                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title = slide.shapes.title
                        title.text = "Scatter Plot"
                        ax = sns.distplot(df[target_column])
                        fig = ax.get_figure()
                        buf = io.BytesIO()
                        fig.savefig(buf, format='png')
                        buf.seek(0)
                        pic = slide.shapes.add_picture(buf, Inches(2), Inches(2), height=Inches(5))
                        top = Inches(7)
                        left = Inches(4.5)
                        width = height = Inches(1)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = "<Project Number> <Project name>\n<Reference number> <Issue> <Last updated>"
                        tf.paragraphs[0].font.size = Pt(8)
                        tf.paragraphs[1].font.size = Pt(8)
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        tf.paragraphs[1].alignment = PP_ALIGN.CENTER

                    if st.checkbox('🈁: Box Plot', value=False):
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title = slide.shapes.title
                        title.text = "Box Plot"
                        top = Inches(7)
                        left = Inches(4.5)
                        width = height = Inches(1)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = "<Project Number> <Project name>\n<Reference number> <Issue> <Last updated>"
                        tf.paragraphs[0].font.size = Pt(8)
                        tf.paragraphs[1].font.size = Pt(8)
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
                        try:
                            if len(box_cols) != 0:
                                for col in box_cols:
                                    try:
                                        data = pd.concat([df[target_column], df[col]], axis=1)
                                        f, ax = plt.subplots(figsize=(8, 6))
                                        fig = sns.boxplot(x=col, y=target_column, data=data)
                                        fig.axis(ymin=np.min(df[target_column]), ymax=np.max(df[target_column]))
                                    except:
                                        st.error('Invalid column')
                                fig = ax.get_figure()
                                buf = io.BytesIO()
                                fig.savefig(buf, format='png')
                                buf.seek(0)
                                pic = slide.shapes.add_picture(buf, Inches(2), Inches(2), height=Inches(5))
                            else:
                                st.error('Box Plot columns not selected')
                                #box_cols = list(df.columns)
                                #if len(box_cols) > 3:
                                #    box_cols = box_cols[:2]
                        except NameError:
                            st.error('Box Plot columns not selected')
                            #box_cols = list(df.columns)
                            #if len(box_cols) > 3:
                            #    box_cols = box_cols[:2]

                    if st.checkbox('🧮: Correlation Matrix', value=False):
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title = slide.shapes.title
                        title.text = "Correlation Matrix"
                        top = Inches(7)
                        left = Inches(4.5)
                        width = height = Inches(1)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = "<Project Number> <Project name>\n<Reference number> <Issue> <Last updated>"
                        tf.paragraphs[0].font.size = Pt(8)
                        tf.paragraphs[1].font.size = Pt(8)
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
                        try:
                            if len(eda_cols) != 0:
                                pass
                            else:
                                eda_cols = list(df.columns)
                                if len(eda_cols) > 5:
                                    eda_cols = eda_cols[:4]
                        except NameError:
                            eda_cols = list(df.columns)
                            if len(eda_cols) > 5:
                                eda_cols = eda_cols[:4]
                        ax = sns.heatmap(df[eda_cols].corr(), cmap="OrRd", linecolor='white', linewidths=1)
                        fig = ax.get_figure()
                        buf = io.BytesIO()
                        fig.savefig(buf, format='png')
                        buf.seek(0)
                        pic = slide.shapes.add_picture(buf, Inches(2), Inches(2), height=Inches(5))

                    if st.checkbox('➿: Pair Plot', value=False):
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title = slide.shapes.title
                        title.text = "Pair Plot"
                        top = Inches(7)
                        left = Inches(4.5)
                        width = height = Inches(1)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.text = "<Project Number> <Project name>\n<Reference number> <Issue> <Last updated>"
                        tf.paragraphs[0].font.size = Pt(8)
                        tf.paragraphs[1].font.size = Pt(8)
                        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
                        try:
                            if len(pair_cols) != 0:
                                pass
                            else:
                                pair_cols = []
                        except NameError:
                            pair_cols = []
                        cols = [target_column]
                        for col in pair_cols:
                            cols.append(col)
                        try:
                            fig = sns.pairplot(df[cols])
                            buf = io.BytesIO()
                            fig.savefig(buf, format='png')
                            buf.seek(0)
                            pic = slide.shapes.add_picture(buf, Inches(2), Inches(2), height=Inches(5))
                        except:
                            st.error('Invalid column')

                    pptx_path = st.text_input(label='How should we call the Presentation?',
                                            value='default_eda.pptx')
                    if st.button("Create PPTX"):
                        prs.save(pptx_path)
                        st.write("Finished! Power Point PPTX was created.")

                if st.sidebar.checkbox('📓 Create Notebook', value=False):
                    st.markdown('### Create a Jupyter Notebook\n'
                                'Which EDA Tools you want to include in the notebook?')

                    nb = nbf.v4.new_notebook()
                    title_text = "# Auto-generated EDA notebook.\n" \
                                "## Welcome to EDA Notebook." \
                                "path of the CSV file.*"
                    import_libraries_code = """# Import all necessary libraries:
    import pandas as pd
    import numpy as np
    import seaborn as sns
    from matplotlib import pyplot as plt"""

                    load_dataframe_code = "# Load the dataframe:\n" \
                                        "# Insert path to CSV:\n" \
                                        """df = pd.read_csv("filename.csv")\ntarget_column = "%s" """ % target_column

                    nb['cells'] = [nbf.v4.new_markdown_cell(title_text),
                                nbf.v4.new_code_cell(import_libraries_code),
                                nbf.v4.new_code_cell(load_dataframe_code)]

                    if st.checkbox('🥩 - Raw Data and Info', value=False):
                        raw_data_code = """# Raw Data and Info:
    print(df.info())
    df.head()"""
                        nb['cells'].append(nbf.v4.new_code_cell(raw_data_code))

                    if st.checkbox('👁 - Missing Data', value=False):
                        missing_data_code = """# Missing Data:
    total = df.isnull().sum().sort_values(ascending=False)
    percent = (df.isnull().sum() / df.isnull().count()).sort_values(ascending=False)
    missing_data = pd.concat([total, percent], axis=1, keys=['Total', 'Percent'])
    missing_data"""
                        nb['cells'].append(nbf.v4.new_code_cell(missing_data_code))

                    if st.checkbox('🔢 - Value Counts', value=False):
                        value_counts_code = """# Value Counts:
    col = df.columns
    # select which column you want to see:
    df[target_column].value_counts()"""
                        nb['cells'].append(nbf.v4.new_code_cell(value_counts_code))

                    if st.checkbox('🧬 - Unique Elements', value=False):
                        unique_elements_code = """# Unique Elements:
    # Unique Elements
    # select which column you want to see:
    df[target_column].unique()"""

                        nb['cells'].append(nbf.v4.new_code_cell(unique_elements_code))

                    if st.checkbox('〽 - Distributions', value=False):
                        distributions_code = """# Distributions:
    print("Skewness: %.3f" % df[target_column].skew())
    print("Kurtosis: %.3f" % df[target_column].kurt())
    sns.distplot(df[target_column])"""
                        nb['cells'].append(nbf.v4.new_code_cell(distributions_code))

                    if st.checkbox('🈁 - Box Plot', value=False):
                        box_plot_code = """# Box Plot:
    # select your variable to make the box plot:
    box_column = df.columns[1]
    data = pd.concat([df[target_column], df[box_column]], axis=1)
    f, ax = plt.subplots(figsize=(8, 6))
    fig = sns.boxplot(x=box_column, y=target_column, data=data)
    fig.axis(ymin=np.min(df[target_column]), ymax=np.max(df[target_column]))"""
                        nb['cells'].append(nbf.v4.new_code_cell(box_plot_code))

                    if st.checkbox('➿ - Pair Plot', value=False):
                        pair_plot_code = """# Pair Plot:
    # select your variable/s to do the pair plot:
    pair_cols = df.columns[1:2]
    pair_cols.append(target_column)
    sns.set()
    sns.pairplot(df[pair_cols], size=2.0)"""
                        nb['cells'].append(nbf.v4.new_code_cell(pair_plot_code))

                    if st.checkbox('🧮 - Correlation Matrix', value=False):
                        correlation_matrix_plot_code = """# Pair Plot:
    # select the columns for the matrix::
    cm_cols = list(df.columns)
    plt.figure(figsize=(12, 8))
    ax = sns.heatmap(df[cm_cols].corr(), cmap="OrRd", linecolor='white', linewidths=1)"""
                        nb['cells'].append(nbf.v4.new_code_cell(correlation_matrix_plot_code))

                    notebook_path = st.text_input(label='How should we call the Jupyter notebook?',
                                                value='default_eda.ipynb')
                    if st.button("Create Notebook"):
                        with open(notebook_path, 'w') as f:
                            nbf.write(nb, f)
                        st.write("Finished! Jupyter Notebook was created.")

    if __name__ == '__main__':
        main()